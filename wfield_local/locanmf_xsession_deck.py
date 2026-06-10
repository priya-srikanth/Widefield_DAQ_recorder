"""Assemble the 6/5-onward CROSS-SESSION deck (Priya 2026-06-10), on freshly re-registered LocaNMF.

Grouped ANIMAL-first, date-second, per LOCANMF_XSESSION_DECK_SPEC.md (sections A-F):
  A  per animal -> date : per-session decode (cue-2s, first-lick-2s, pre-cue-2s) + per-animal rolling
  B  per date -> all animals : cue-aligned rolling dynamics + 2s-vs-rolling
  C  per animal : per-position decoder weights + region ablation (leave-1/leave-2) across sessions
  D  per date -> all animals : encoder time-local quiet baseline
  E  per animal : expected MO/SS activity by position + encoding EV per region (abs + FEVE 1.0)
  F  cross-mouse representation + within-animal consistency + RSA/RDM (6/5-6/8)

Consumes PNGs produced by build_xsession_figs.py. Missing images are skipped.

    python -c "from pathlib import Path; from wfield_local.locanmf_xsession_deck import build_xsession_ppt; \
               build_xsession_ppt(Path('C:/Users/sabatini/source/cue_lick'))"
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

NAVY = RGBColor(0x1F, 0x33, 0x55); GREY = RGBColor(0x55, 0x55, 0x55); RED = RGBColor(0xB0, 0x22, 0x22)
DATES = [("0605", "6/5"), ("0606", "6/6"), ("0607", "6/7"), ("0608", "6/8")]
ANIMALS = ["PS92", "PS93", "PS94", "PS95"]
TAG = "0605-0608"


def build_xsession_ppt(src: Path, out_name="spout_position_decoder_xsession_6on.pptx") -> Path:
    src = Path(src)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def slide():
        return prs.slides.add_slide(BLANK)

    def title(s, text, sub=None, subcolor=GREY):
        tf = s.shapes.add_textbox(Inches(0.4), Inches(0.16), Inches(12.6), Inches(0.95)).text_frame
        tf.word_wrap = True; r = tf.paragraphs[0].add_run(); r.text = text
        r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = NAVY
        if sub:
            r2 = tf.add_paragraph().add_run(); r2.text = sub; r2.font.size = Pt(12.5); r2.font.color.rgb = subcolor

    def pic(s, p, **kw):
        p = Path(p)
        if p.exists():
            s.shapes.add_picture(str(p), **kw); return True
        return False

    def big(s, p, top=1.35, width=12.7):
        w = Inches(width); pic(s, p, left=(SW - w) / 2, top=Inches(top), width=w)

    def grid(s, paths, cols=2, top=1.25, side=0.25, gap=0.18, bottom=0.25):
        paths = [Path(p) for p in paths]
        if not any(p.exists() for p in paths):
            return
        rows = (len(paths) + cols - 1) // cols
        cell_w = (SW - Inches(side) * 2 - Inches(gap) * (cols - 1)) / cols
        cell_h = (SH - Inches(top) - Inches(bottom) - Inches(gap) * (rows - 1)) / rows
        for i, p in enumerate(paths):
            if not p.exists():
                continue
            r, c = divmod(i, cols)
            iw, ih = Image.open(str(p)).size           # fit within BOTH cell dims, keep aspect
            scale = min(cell_w / iw, cell_h / ih)
            w, h = int(iw * scale), int(ih * scale)
            left = Inches(side) + c * (cell_w + Inches(gap)) + (cell_w - w) / 2
            t = Inches(top) + r * (cell_h + Inches(gap)) + (cell_h - h) / 2
            s.shapes.add_picture(str(p), left, t, width=w, height=h)

    def divider(text, sub=None):
        s = slide()
        tf = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.8)).text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = text; r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
        if sub:
            r2 = tf.add_paragraph().add_run(); r2.text = sub; r2.font.size = Pt(15); r2.font.color.rgb = GREY

    def sess(label, align):
        return src / f"locanmf_position_session_{label}_locanmf_{align}_base-none_cv-block.png"

    # ---------------- title ----------------
    s = slide()
    tf = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.0)).text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "Spout-position decoding & encoding from cortex"
    r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = NAVY
    for t in ["Cross-session re-registration (each session aligned to that animal's 6/6 image) — baseline 6/5–6/8",
              "PS92 / PS93 / PS94 / PS95   (PS93 = right orofacial deficit)",
              "Grouped by animal, then date. Individual LocaNMF components, block-aware CV, chance = 0.17"]:
        rr = tf.add_paragraph().add_run(); rr.text = t; rr.font.size = Pt(15); rr.font.color.rgb = GREY

    # ---------------- Section A ----------------
    divider("A. Per-animal decoding across sessions",
            "Cue-anchored 2 s (post-cue) and first-lick 2 s confusion + recall; rolling pre-cue→post-cue; pre-cue 2 s maintained code.")
    ALN = [("lick", "first-lick 2 s (headline engaged decoder)"),
           ("cue", "cue-aligned 2 s post-cue (engaged blue vs no-lick red)"),
           ("precue", "pre-cue 2 s (engaged + no-lick = maintained position code)")]
    for a in ANIMALS:
        for align, desc in ALN:
            s = slide()
            sub = "Per session: confusion matrix + per-position recall."
            if a == "PS92" and align == "precue":
                sub += "  NB PS92 6/5 ENL was 1.5–2.5 s — the 2 s pre-cue window may clip a few licks (others 2–3 s)."
            title(s, f"{a} — {desc}", sub)
            grid(s, [sess(f"{a}_{d}", align) for d, _ in DATES], cols=2)
        s = slide()
        title(s, f"{a} — rolling decoder across sessions (pre-cue ENL → post-cue)",
              "Sliding 0.5 s window, block-CV, one line per session. Ramp through ENL = anticipatory; early-ENL level above chance = maintained code.")
        big(s, src / f"locanmf_decoder_rolling_by_animal_{a}.png", top=1.5, width=11.2)

    # ---------------- Section B ----------------
    divider("B. Per-date rolling dynamics (all animals)",
            "Cue-aligned rolling decoder + does the temporal profile beat the 2 s window-mean?")
    for d, dl in DATES:
        s = slide()
        title(s, f"{dl} — cue-aligned rolling decoder (all animals)",
              "Position decodability through the pre-cue ENL → post-cue. One line per animal, sliding 0.5 s window, block-CV.")
        big(s, src / f"locanmf_decoder_rolling_cue_{d}.png", top=1.5, width=11.0)
        s = slide()
        title(s, f"{dl} — rolling temporal bins vs single 2 s window",
              "First-lick aligned. Left: accuracy vs time. Right: 2 s mean vs temporal bins — does the time-course add information?")
        big(s, src / f"locanmf_decoder_temporal_dynamics_{d}.png", top=1.45, width=12.6)

    # ---------------- Section C ----------------
    divider("C. Per-animal decoder weights & region ablation",
            "Which regions code each position, and region necessity (leave-one vs leave-two-out) across that animal's sessions.")
    for a in ANIMALS:
        s = slide()
        title(s, f"{a} — per-position decoder weight by Allen region (across sessions)",
              "Positive LR weight per region, one panel per session. Contralateral SSp: left spouts→right-hemi, right→left.")
        big(s, src / f"locanmf_decoder_weights_by_region_{a}.png", top=1.5, width=13.0)
        s = slide()
        title(s, f"{a} — region ablation: leave-one vs leave-two-out (across sessions)",
              "Unilateral sensory/motor. Leave-TWO exposes necessity hidden by cross-hemisphere redundancy. One panel per session.")
        big(s, src / f"locanmf_ablation_grouped_unilateral_{a}.png", top=1.5, width=13.0)

    # ---------------- Section D ----------------
    divider("D. Encoder baseline — time-local quiet reference",
            "Pooled rest baseline drift over each session; the time-local zero for the pre/post-stroke residual.")
    for d, dl in DATES:
        s = slide()
        title(s, f"{dl} — encoder time-local quiet (rest) baseline (all animals)",
              "Drift of the pooled rest baseline across the session (ΔF/F). The local baseline tracks it → stable zero.")
        big(s, src / f"locanmf_encoder_quiet_drift_{d}.png", top=1.5, width=10.6)

    # ---------------- Section E ----------------
    divider("E. Per-animal encoder: expected activity & explained variance",
            "Position → expected cortical activity (SSp / MO) and encoding explained variance per region (absolute + normalized to 1.0 = FEVE).")
    for a in ANIMALS:
        s = slide()
        title(s, f"{a} — expected SSp / MO activity by position (encoder, per session)",
              "Predicted per-position time-course of pooled SSp and MO activity (lick-aligned). One panel per session.")
        grid(s, [src / f"locanmf_encoder_temporal_{a}_{d}.png" for d, _ in DATES], cols=2)
        s = slide()
        title(s, f"{a} — encoding explained variance per region: absolute & normalized (FEVE)",
              "Left of each: absolute explainable-vs-captured; right: FEVE (captured / explainable, →1.0). One panel per session.")
        grid(s, [src / f"locanmf_encoder_r2_by_region_{a}_{d}.png" for d, _ in DATES], cols=2)
    s = slide()
    title(s, "Encoder — FEVE by region, pooled per animal (6/5–6/8)",
          "Fraction of EXPLAINABLE variance captured per region, pooled over each animal's re-registered sessions.")
    big(s, src / "locanmf_encoder_feve_by_region_pooled.png", top=1.5, width=12.9)
    s = slide()
    title(s, "Encoder — FEVE by region, individual sessions per animal (6/5–6/8)",
          "Same metric, one row per session → session-to-session stability. PS93's SSp-left vs -right split is visible.")
    big(s, src / "locanmf_encoder_feve_by_region_sessions.png", top=1.5, width=12.9)

    # ---------------- Section F ----------------
    divider("F. Cross-mouse representation, consistency & geometry (6/5–6/8)")
    s = slide()
    title(s, "Cross-mouse cortical representation of spout position (6/5–6/8)",
          "Overall/per-position decoding, L-vs-R spout decodability, SSp-left-vs-right hemisphere, encoding EV, L/R asymmetry. "
          "PS93 right orofacial deficit → predicts SSp-left ≪ -right.")
    big(s, src / f"locanmf_cross_mouse_comparison_{TAG}.png", top=1.5, width=12.7)
    s = slide()
    title(s, "Within-animal consistency of per-position decode / encode (6/5–6/8)",
          "Per-position profile per session (grey) + mean ± SD (bold). Mean SD = the noise floor a post-stroke change must clear.")
    big(s, src / f"locanmf_within_animal_consistency_{TAG}.png", top=1.5, width=12.9)
    s = slide()
    title(s, "RSA — representational geometry of spout position (6/5–6/8)",
          "6×6 position RDM per session; 2nd-order RSA (basis-free). Within-animal > across-animal = stable individual geometry; "
          "% = within / split-half noise ceiling.")
    big(s, src / f"locanmf_rsa_sessions_{TAG}.png", top=1.6, width=13.0)
    s = slide()
    title(s, "RSA — mean representational dissimilarity matrix per animal (6/5–6/8)",
          "How the 6 positions relate (dark = similar patterns, bright = distinct).")
    big(s, src / f"locanmf_rsa_rdms_{TAG}.png", top=1.9, width=12.7)
    s = slide()
    title(s, "Hemisphere-resolved RDM — PS93's two hemispheres encode position differently (6/5–6/8)",
          "Disattenuated left-vs-right RDM agreement; PS93 lowest. Its LEFT hemisphere is reliable but its geometry is reshaped, "
          "not abolished — the lateralized signature the pooled RDM misses.")
    big(s, src / f"locanmf_rsa_hemisphere_summary_{TAG}.png", top=1.6, width=13.0)
    s = slide()
    title(s, "Hemisphere-resolved RDM — left-hem (top) vs right-hem (bottom) geometry per animal (6/5–6/8)")
    big(s, src / f"locanmf_rsa_hemisphere_rdms_{TAG}.png", top=1.9, width=13.0)

    outp = src / out_name; prs.save(str(outp)); return outp


if __name__ == "__main__":
    import sys
    out = build_xsession_ppt(Path(sys.argv[1] if len(sys.argv) > 1 else "C:/Users/sabatini/source/cue_lick"))
    print("wrote", out)
