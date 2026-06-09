"""Assemble the spout-position decoder summary deck from figures in a directory. Expects the
per-day decoder figures (from locanmf_position_decoder.py, --align lick/cue/precue) and the
analysis figures (from locanmf_decoder_weights.py) to be present; missing images are skipped."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1F, 0x33, 0x55); GREY = RGBColor(0x55, 0x55, 0x55); RED = RGBColor(0xB0, 0x22, 0x22)
DAYS = [("0601", "6/1"), ("0602", "6/2"), ("0603", "6/3"), ("0604", "6/4"), ("0605", "6/5"), ("0606", "6/6"), ("0607", "6/7")]


def build_ppt(src: Path, out_name="spout_position_decoder_summary.pptx") -> Path:
    src = Path(src)

    def dfig(day, align, kind):
        return src / f"locanmf_position_{kind}_{day}_locanmf_{align}_base-none_cv-block.png"

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def title(slide, text, sub=None):
        tf = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.9)).text_frame
        tf.word_wrap = True; r = tf.paragraphs[0].add_run(); r.text = text
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
        if sub:
            r2 = tf.add_paragraph().add_run(); r2.text = sub; r2.font.size = Pt(13); r2.font.color.rgb = GREY

    def imgs(slide, paths, top=Inches(1.3)):
        paths = [p for p in paths if p.exists()]
        if not paths:
            return
        w = Inches(6.2); gap = Inches(0.25); total = len(paths) * w + (len(paths) - 1) * gap
        left0 = (prs.slide_width - total) / 2
        for i, p in enumerate(paths):
            slide.shapes.add_picture(str(p), left0 + i * (w + gap), top, width=w)

    def pic(slide, p, **kw):
        if p.exists():
            slide.shapes.add_picture(str(p), **kw)

    # title
    s = prs.slides.add_slide(BLANK)
    tf = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.8)).text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "Spout-position decoding from cortex"
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NAVY
    for t in ["Best strategy: individual LocaNMF components, NO per-trial baseline, block-aware CV, first-lick 2 s",
              "PS92 / PS94 / PS95  —  baseline days 6/1-6/4 (all pre-stroke)",
              "Multinomial logistic regression - 6 positions - chance = 0.17"]:
        rr = tf.add_paragraph().add_run(); rr.text = t; rr.font.size = Pt(16); rr.font.color.rgb = GREY

    # strategy + table
    s = prs.slides.add_slide(BLANK)
    title(s, "Why this strategy", "What changed from the first pass, and the resulting accuracy")
    bt = s.shapes.add_textbox(Inches(0.45), Inches(1.2), Inches(6.0), Inches(5.8)).text_frame; bt.word_wrap = True
    bullets = [
        ("Individual components, not region-pooled", "all LocaNMF components are features; labels only subset SSp/MO."),
        ("No per-trial baseline", "a session-constant (quiet) baseline is invisible to a standardized decoder; per-trial pre-cue over-subtracts real anticipatory signal."),
        ("Block-aware CV", "positions come in ~6-trial blocks; random k-fold leaks each block's slow-drift fingerprint."),
        ("First-lick, 2 s window", "integrates the lick bout (beats 1 s); >~2.5 s dilutes. Cue/pre-cue kept for the no-lick test."),
        ("SSp carries it, MO secondary", "contralateral orofacial SSp dominates; MOp/MOs above chance, strongest for far."),
        ("No-lick (post-cue) ~ chance; pre-cue no-lick above chance", "maintained position code readable without a lick."),
    ]
    first = True
    for h, d in bullets:
        p = bt.paragraphs[0] if first else bt.add_paragraph(); first = False
        rr = p.add_run(); rr.text = "- " + h; rr.font.size = Pt(13.5); rr.font.bold = True; rr.font.color.rgb = NAVY
        p2 = bt.add_paragraph(); r2 = p2.add_run(); r2.text = "   " + d; r2.font.size = Pt(11); r2.font.color.rgb = GREY
    rows = [("", "PS92", "PS94", "PS95"), ("6/3 first-lick", "0.67", "0.83", "0.85"),
            ("6/4 first-lick", "0.56", "0.29", "0.49"), ("6/1 first-lick", "-", "0.73", "0.91"),
            ("6/2 first-lick", "0.46", "-", "-"), ("6/3 pre-cue no-lick", "0.27", "0.34", "0.22")]
    tbl = s.shapes.add_table(len(rows), 4, Inches(6.8), Inches(1.5), Inches(6.0), Inches(3.0)).table
    for ci in range(4):
        tbl.columns[ci].width = Inches(1.5)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = val if val else " "
            par = cell.text_frame.paragraphs[0]; par.alignment = PP_ALIGN.CENTER; run = par.runs[0]; run.font.size = Pt(12)
            if ri == 0 or ci == 0:
                run.font.bold = True; run.font.color.rgb = NAVY
            if "no-lick" in row[0] and ci > 0:
                run.font.color.rgb = RED

    # figure slides
    for day, dlab in DAYS:
        s = prs.slides.add_slide(BLANK)
        title(s, f"{dlab} - first-lick 2 s (best engaged decoder)", "Left: confusion matrix.  Right: per-position recall.")
        imgs(s, [dfig(day, "lick", "decoder"), dfig(day, "lick", "recall")])
    for day, dlab in DAYS:
        s = prs.slides.add_slide(BLANK)
        title(s, f"{dlab} - cue 2 s (no-lick generalization test)", "Right: engaged (blue) vs no-lick (red).")
        imgs(s, [dfig(day, "cue", "decoder"), dfig(day, "cue", "recall")])
    for day, dlab in DAYS:
        s = prs.slides.add_slide(BLANK)
        title(s, f"{dlab} - PRE-CUE 1 s (motor-independent; applies to no-lick)", "No-lick above chance = maintained code.")
        imgs(s, [dfig(day, "precue", "decoder"), dfig(day, "precue", "recall")])

    # analysis slides
    s = prs.slides.add_slide(BLANK)
    title(s, "Which components carry each position - and where", "LR weight by Allen region (top-12 + MOp/MOs L/R).")
    pic(s, src / "locanmf_decoder_weights_by_region.png", left=Inches(0.3), top=Inches(1.35), width=Inches(12.7))
    s = prs.slides.add_slide(BLANK)
    title(s, "SSp vs MO", "Accuracy by feature set + share of decoder weight by region group.")
    pic(s, src / "locanmf_decoder_region_groups.png", left=Inches(0.4), top=Inches(1.5), width=Inches(12.5))
    for lab in ("PS94_0603", "PS95_0605"):
        fp = src / f"locanmf_region_ablation_{lab}.png"
        if fp.exists():
            s = prs.slides.add_slide(BLANK)
            title(s, f"Region importance by ABLATION — {lab}",
                  "region-only = sufficiency; leave-out drop = necessity. SSp leads by sufficiency; low leave-out drops "
                  "= position info is redundant across cortex (PS94 6/3 SSp-localized; PS95 6/5 fully redundant/global -> movement?).")
            pic(s, fp, left=Inches(0.3), top=Inches(1.55), width=Inches(12.7))
    fp = src / "locanmf_ablation_grouped_unilateral.png"
    if fp.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "Unilateral sensory/motor: leave-one vs leave-two-out",
              "Leave-TWO exposes necessity hidden by cross-hemisphere redundancy: PS94 6/3 bilateral sensory IS necessary "
              "(0.044/0.065 one-hemi -> 0.108 both); motor not. PS95 6/5 fully redundant (even all-sensory ~0 -> global/movement).")
        pic(s, fp, left=Inches(0.4), top=Inches(1.6), width=Inches(12.5))
    s = prs.slides.add_slide(BLANK)
    title(s, "Why a 2 s window", "Decoding vs post-cue window length.")
    pic(s, src / "locanmf_decoder_window_sweep.png", left=Inches(2.9), top=Inches(1.35), height=Inches(5.3))
    s = prs.slides.add_slide(BLANK)
    title(s, "Pre-stroke baseline variability (3 days/animal)", "The spread any post-stroke effect must clear.")
    pic(s, src / "locanmf_decoder_baseline_variability.png", left=Inches(0.3), top=Inches(1.4), width=Inches(12.7))
    for day, dlab in DAYS:
        fp = src / f"locanmf_decoder_temporal_dynamics_{day}.png"
        if not fp.exists():
            continue
        s = prs.slides.add_slide(BLANK)
        title(s, f"Rolling temporal dynamics - {dlab}",
              "When is position info present, and does the temporal profile beat the window-mean?")
        pic(s, fp, left=Inches(0.4), top=Inches(1.5), width=Inches(12.5))

    rc = src / "locanmf_decoder_rolling_cue_0607.png"
    if rc.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "6/7 cue-aligned rolling decoder (pre-cue ENL -> post-cue)",
              "Position decodability ramps through the longer 2-3s ENL (maintained code), peaks ~0.5-0.8s post-cue. "
              "PS93 now 3 sessions. Sliding 0.5s window, block-CV.")
        pic(s, rc, left=Inches(1.4), top=Inches(1.5), width=Inches(10.5))
    lr = src / "locanmf_decoder_rolling_laterality_0607.png"
    if lr.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "6/7 rolling decoder: laterality (L/center/R) vs full 6-way",
              "Collapsing to laterality (3-way, chance 0.33) gives higher ABSOLUTE accuracy but similar/slightly WORSE "
              "above-chance than 6-way -- the close/far distinction carries real information; collapsing doesn't help.")
        pic(s, lr, left=Inches(2.0), top=Inches(1.5), width=Inches(9.5))

    s = prs.slides.add_slide(BLANK)
    title(s, "6/4 engagement: first-40-min decoding",
          "Late-session disengagement drags 6/4 down; restricting to the engaged window partly recovers it (not 6/3).")
    pic(s, src / "locanmf_decoder_first40min_0604.png", left=Inches(0.4), top=Inches(1.5), width=Inches(12.5))
    s = prs.slides.add_slide(BLANK)
    title(s, "6/4 v1 vs v2 Allen-CCF alignment",
          "ROI decoder (alignment-sensitive) is unchanged by v2 re-registration -> the 6/4 deficit is not an alignment artifact.")
    pic(s, src / "locanmf_decoder_v1_vs_v2_alignment_0604.png", left=Inches(0.4), top=Inches(1.5), width=Inches(12.5))
    s = prs.slides.add_slide(BLANK)
    title(s, "Pre-cue (ENL) signal + lick timing -> ENL length",
          "Lick-free window is only ~1.4s now; the maintained position code is real but ramps toward the cue. "
          "A longer ENL (2-3s) would give a cleaner motor-independent pre-cue readout.")
    pic(s, src / "locanmf_precue_ENL_analysis_0603.png", left=Inches(0.4), top=Inches(1.6), width=Inches(12.5))

    # top predictive component footprints, per session
    for day, dlab in DAYS:
        for lab in [f"{an}_{day}" for an in ("PS92", "PS94", "PS95")]:
            fp = src / f"locanmf_top_components_{lab}.png"
            if not fp.exists():
                continue
            s = prs.slides.add_slide(BLANK)
            title(s, f"{lab[:4]} {dlab} - top-10 components by univariate decoding accuracy",
                  "Ranked by per-component block-CV accuracy (NOT |weight|, which surfaces suppressors). Allen-overlaid footprints.")
            pic(s, fp, left=Inches(0.3), top=Inches(1.6), width=Inches(12.7))

    # ENCODER section (reverse model: position -> expected activity)
    encs = [f"{an}_0607" for an in ("PS92", "PS93", "PS94", "PS95")]
    if any((src / f"locanmf_encoder_predicted_maps_{l}.png").exists() for l in encs):
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — position -> expected neural activity (reverse of the decoder)",
              "Fit pre-stroke; post-stroke the residual (observed - predicted) per INTENDED position = the lesion's "
              "effect, computable even on no-lick/failed trials. Predicted maps + encoding R^2 + expected dynamics follow.")
    vs = src / "locanmf_encoder_vs_svd_PS95_0607.png"
    if vs.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "VALIDATION: encoder (LocaNMF) maps vs SVD pixel maps",
              "Matched cue-aligned pre-cue-delta. Per-position spatial r = 0.99-1.00 (all sessions) -> LocaNMF reconstruction "
              "is the same as the SVD pixel data; the encoder's predicted maps are the genuine cortical patterns.")
        pic(s, vs, left=Inches(0.3), top=Inches(2.2), width=Inches(12.7))
    evp = src / "locanmf_encoder_ev_by_position_0607.png"
    if evp.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — explained variance per spout position (per session)",
              "Held-out R^2 restricted to each position's trials. Lateral/extreme positions are most distinctly encoded; "
              "center positions sit near the grand mean (low/negative). PS94 weak throughout.")
        pic(s, evp, left=Inches(0.8), top=Inches(1.7), width=Inches(11.5))
    evc = src / "locanmf_encoder_ev_ceiling_by_position_0607.png"
    if evc.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — explained variance per position RELATIVE TO CEILING",
              "Left: explainable variance (noise ceiling) per position. Right: captured/ceiling. Center positions have "
              "~0 ceiling (no position-distinct signal) -> their low raw EV is nothing-to-explain, not encoder failure.")
        pic(s, evc, left=Inches(0.3), top=Inches(1.7), width=Inches(12.7))
    fvp = src / "locanmf_encoder_feve_by_region_pooled.png"
    if fvp.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — FEVE by region, pooled per animal (across ALL sessions)",
              "FEVE = % of EXPLAINABLE (ceiling) variance the encoder captures, per region (SS pooled over each animal's "
              "sessions). 100% = all position-explainable variance captured; SSp/MO labels in red/blue. Region axis = "
              "regions with non-trivial explainable signal, sorted by explainable variance.")
        pic(s, fvp, left=Inches(0.2), top=Inches(1.6), width=Inches(12.9))
    fvs = src / "locanmf_encoder_feve_by_region_sessions.png"
    if fvs.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — FEVE by region, individual sessions per animal",
              "Same ceiling-normalized metric, one row per animal with a bar per session (date) -> session-to-session "
              "stability of each region's encoding. PS93's SSp-left vs -right split is visible here too.")
        pic(s, fvs, left=Inches(0.2), top=Inches(1.5), width=Inches(12.9))
    qd = src / "locanmf_encoder_quiet_drift_0607.png"
    if qd.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER baseline: time-local quiet (rest) reference",
              "Pooled rest baseline over the session. Drift is small in dF/F (+-0.005; PS92 a clear downward trend, PS93 "
              "jumpy over 123 min). The time-local baseline tracks it -> stable zero for the pre/post-stroke residual.")
        pic(s, qd, left=Inches(1.6), top=Inches(1.7), width=Inches(10.0))
    for lab in encs:
        for kind, sub in [("predicted_maps", "expected cortical activity per intended position (footprint-reconstructed)"),
                          ("temporal", "expected activity time-course per position (SSp / MO pooled, lick-aligned)"),
                          ("r2_by_region", "cross-validated encoding R^2 by region (activity explained by position)")]:
            fp = src / f"locanmf_encoder_{kind}_{lab}.png"
            if fp.exists():
                s = prs.slides.add_slide(BLANK)
                title(s, f"ENCODER {lab} — {kind.replace('_', ' ')}", sub)
                w = Inches(12.5) if kind != "r2_by_region" else Inches(7.5)
                pic(s, fp, left=(prs.slide_width - w) / 2, top=Inches(1.55), width=w)

    # cross-mouse comparison
    cm = src / "locanmf_cross_mouse_comparison.png"
    if cm.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "CROSS-MOUSE: cortical representation of spout position (all sessions)",
              "Per mouse: overall/per-position decoding, left-vs-right spout decodability, SSp-left-vs-right hemisphere, "
              "per-position encoding EV, and L/R asymmetry indices. PS93 has a RIGHT orofacial deficit (predicts L/R asymmetry).")
        pic(s, cm, left=Inches(0.3), top=Inches(1.5), width=Inches(12.7))
    wac = src / "locanmf_within_animal_consistency.png"
    if wac.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "WITHIN-ANIMAL consistency of per-position decode/encode across sessions",
              "Each animal's per-position profile per session (grey) + mean +- SD (bold). Pairwise r = pattern "
              "reproducibility, mean SD = magnitude noise floor a post-stroke change must clear. Decode SD ~0.15-0.19 "
              "(6-session animals span the noisy/low-engagement early days); consistency is engagement-dependent.")
        pic(s, wac, left=Inches(0.2), top=Inches(1.5), width=Inches(12.9))
    wac3 = src / "locanmf_within_animal_consistency_0605-0607.png"
    if wac3.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "WITHIN-ANIMAL consistency on 3 consecutive matched-engagement days (6/5-6/7)",
              "Same analysis restricted to 6/5-6/7. Decode per-position SD tightens to ~0.04-0.08 (vs ~0.15-0.19 "
              "across all baseline days) -> the all-sessions variability was the noisy/low-engagement early days "
              "(6/1,6/4), not intrinsic. NB pairwise r understates consistency for flat high-recall profiles (PS95 "
              "r=0.30 but SD=0.04) -> trust SD there. Engagement-matched consecutive baselines => noise floor ~0.05.")
        pic(s, wac3, left=Inches(0.2), top=Inches(1.5), width=Inches(12.9))
    rsa = src / "locanmf_rsa_sessions.png"
    if rsa.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "RSA — representational geometry of spout position (within vs across animals)",
              "Per session a 6x6 RDM (1 - corr between the 6 position activity patterns); 2nd-order RSA = Spearman "
              "between RDMs (basis-free, valid across sessions/animals). Within-animal RDM similarity > across-animal "
              "for all 4 (stable individual geometry); % = within / split-half noise ceiling. PS93 is NOT the geometric "
              "outlier (most similar to PS92) -> its deficit is the lateralized SSp asymmetry (F15), not a global "
              "geometry change. PS94 most distinct.")
        pic(s, rsa, left=Inches(0.15), top=Inches(1.7), width=Inches(13.0))
    rsr = src / "locanmf_rsa_rdms.png"
    if rsr.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "RSA — mean representational dissimilarity matrix per animal",
              "How the 6 positions relate (dark = similar patterns, bright = distinct). PS93 flattest (positions less "
              "differentiated, but rank-ordered like PS92); PS95 far_center the standout-distinct position (reproduces F7).")
        pic(s, rsr, left=Inches(0.3), top=Inches(2.2), width=Inches(12.7))

    # takeaways
    s = prs.slides.add_slide(BLANK); title(s, "Takeaways")
    tf = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.6)).text_frame; tf.word_wrap = True
    takes = [
        "Intended spout position decodes well above chance on every baseline session (first-lick 2 s, block-CV; up to 0.91).",
        "SSp orofacial subfields, contralateral to the spout, dominate; MOp/MOs are secondary (strongest for far positions).",
        "Post-cue no-lick decodes at chance (negative control); PRE-CUE no-lick decodes above chance = motor-independent maintained code.",
        "Baseline day-to-day variability is large (PS94 0.29-0.83); 6/4 is the low/low-engagement day.",
        "For the stroke pre/post: block-CV, no per-trial baseline, multiple baseline days, engagement matching, and the pre-cue readout.",
    ]
    first = True
    for t in takes:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = "- " + t; r.font.size = Pt(15); r.font.color.rgb = NAVY; p.space_after = Pt(10)

    outp = src / out_name; prs.save(str(outp)); return outp
