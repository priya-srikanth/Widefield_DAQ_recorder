"""Add a PS92 section to the PS94/PS95 alignment comparison deck.

Copies the latest deck to E:\\labcams_data\\PS92_94_95.pptx and appends a PS92
section matching the existing per-animal slide/sizing conventions:
  divider -> mean+Allen -> spout avg -> spout shared-scale -> spout pairwise ->
  post-lick maps -> post-lick pairwise -> cue-vs-lick.
"""
import shutil, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

SRC = r"E:\labcams_data\20260601\alignment_comparison_PS94_PS95_PS.pptx"
DST = r"E:\labcams_data\PS92_94_95.pptx"
shutil.copy2(SRC, DST)
prs = Presentation(DST)
blank = next(l for l in prs.slide_layouts if l.name == "Blank")

MC = r"E:\labcams_data\20260602\PS92\PS92_20260602_151820\illuminated_rescue\motion_corrected"
SP = os.path.join(MC, "spout_trial_averages_allen_v2")
LK = os.path.join(MC, "lick_aligned_v2")


def txt(slide, l, t, w, h, s, size, color, bold=False, align=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = s
    if align is not None:
        p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)
    return tb


def set_dark_bg(slide, hexcolor="222222"):
    cSld = slide._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
    bg = etree.SubElement(cSld, "{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
    bgPr = etree.SubElement(bg, "{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr")
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sf = etree.SubElement(bgPr, f"{{{a}}}solidFill")
    c = etree.SubElement(sf, f"{{{a}}}srgbClr"); c.set("val", hexcolor)
    etree.SubElement(bgPr, f"{{{a}}}effectLst")
    cSld.insert(0, bg)  # bg must come before spTree


def content(title, subtitle, img, l, t, w, footer):
    s = prs.slides.add_slide(blank)
    txt(s, 0.44, 0.25, 12.44, 0.32, title, 19, "222222", bold=True)
    txt(s, 0.44, 0.64, 12.44, 0.30, subtitle, 8.5, "666666")
    if img and os.path.exists(img):
        s.shapes.add_picture(img, Inches(l), Inches(t), width=Inches(w))
    txt(s, 0.44, 7.19, 12.44, 0.14, footer, 6.5, "666666")
    return s


# update title slide
for sh in prs.slides[0].shapes:
    if sh.has_text_frame and "PS94 and PS95" in sh.text_frame.text:
        sh.text_frame.paragraphs[0].runs[0].text = "PS92, PS94, and PS95 Allen Alignment Comparison"

# ---- PS92 divider ----
d = prs.slides.add_slide(blank)
set_dark_bg(d, "222222")
txt(d, 0.64, 2.85, 11.81, 0.74, "PS92", 44, "FFFFFF", bold=True)
txt(d, 0.69, 3.86, 11.53, 0.4,
    "Rescued recording (illuminated-frame consolidation). Functional-channel fix + ROI-aware Allen alignment (reference space).",
    14, "D8D8D8")

# ---- PS92 content slides (match PS95 sizing) ----
content("PS92 mean 415/470 nm with Allen outlines",
        "Rescued, hemodynamic-corrected (functional channel corrected); ROI-aware Allen alignment warped to atlas grid (reference space).",
        os.path.join(SP, "PS92_v2_mean_415_470_with_allen_overlay.png"), 0.10, 1.95, 13.14,
        os.path.join(SP, "PS92_v2_mean_415_470_with_allen_overlay.png"))

content("PS92 spout-position averages",
        "1 s pre-cue, 1 s post-cue, and post-pre delta maps by spout position",
        os.path.join(SP, "PS92_v2_spout_positions_1s_pre_post_delta_allen_overlay.png"), 4.91, 1.59, 3.51,
        os.path.join(SP, "PS92_v2_spout_positions_1s_pre_post_delta_allen_overlay.png"))

content("PS92 spout-position averages - shared scale",
        "Same diverging color scale across all pre, post, and post-pre panels in this figure",
        os.path.join(SP, "PS92_v2_spout_positions_1s_pre_post_delta_shared_scale.png"), 4.91, 1.59, 3.51,
        os.path.join(SP, "PS92_v2_spout_positions_1s_pre_post_delta_shared_scale.png"))

content("PS92 pairwise delta-position contrasts",
        "Each panel is first condition's post-pre map minus second condition's post-pre map",
        os.path.join(SP, "PS92_v2_pairwise_spout_position_delta_contrasts_allen_overlay.png"), 4.01, 1.59, 5.31,
        os.path.join(SP, "PS92_v2_pairwise_spout_position_delta_contrasts_allen_overlay.png"))

content("PS92 post-lick maps by spout position",
        "Mean 150 ms post-lick by spout position (post-lick window; see notes on bout overlap)",
        os.path.join(LK, "PS92_v2_lick_aligned_150ms_post_by_spout.png"), 2.60, 1.40, 8.20,
        os.path.join(LK, "PS92_v2_lick_aligned_150ms_post_by_spout.png"))

content("PS92 post-lick pairwise spout-position contrasts",
        "Each panel is first position's post-lick map minus second position's post-lick map",
        os.path.join(LK, "PS92_v2_lick_aligned_pairwise_spout_position_contrasts.png"), 4.01, 1.59, 5.31,
        os.path.join(LK, "PS92_v2_lick_aligned_pairwise_spout_position_contrasts.png"))

content("PS92 cue-aligned versus lick-aligned maps",
        "Post-cue (1 s) vs post-lick (150 ms) mean maps by spout position",
        os.path.join(LK, "PS92_v2_cue_vs_lick_spout_position_maps.png"), 5.00, 1.59, 3.20,
        os.path.join(LK, "PS92_v2_cue_vs_lick_spout_position_maps.png"))

prs.save(DST)
print("saved", DST, "num slides:", len(prs.slides))
