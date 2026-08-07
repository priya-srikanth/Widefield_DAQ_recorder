from pptx import Presentation
from pptx.util import Emu
import sys

src = r"E:\labcams_data\20260601\alignment_comparison_PS94_PS95_PS.pptx"
prs = Presentation(src)
EMU = 914400.0
print("slide size: %.2f x %.2f in" % (prs.slide_width / EMU, prs.slide_height / EMU))
print("num slides:", len(prs.slides))
for si, slide in enumerate(prs.slides):
    print(f"\n===== SLIDE {si} (layout={slide.slide_layout.name}) =====")
    for sh in slide.shapes:
        t = sh.shape_type
        pos = "L=%.2f T=%.2f W=%.2f H=%.2f" % (
            (sh.left or 0) / EMU, (sh.top or 0) / EMU, (sh.width or 0) / EMU, (sh.height or 0) / EMU)
        txt = ""
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = " | TEXT: " + repr(sh.text_frame.text.strip()[:90])
        img = ""
        if t == 13:  # PICTURE
            try:
                img = " | IMG name=%s" % (sh.name,)
            except Exception:
                pass
        print(f"  [{si}.{sh.shape_id}] {t} {pos}{img}{txt}")
