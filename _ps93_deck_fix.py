"""Replace stale PS93 8/5 '4/6 INCOMPLETE' red caveat boxes with a resolved recovery note (green)."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
DECK=r"N:\MICROSCOPE\Priya\Widefield\labcams\PS92_94_95_affine8v1.pptx"
NEW=("✓ All 6 spout positions RECOVERED from cam1 head-on video. DAQ strobe bit1 was dead "
     "8/5 AND this session's behavior log is empty, so positions were inferred from the spout's "
     "x-position in each DAQ-sync-aligned cam1 frame (discrete clusters per fixed detent; min "
     "margin 10 px from decision threshold; counts match balanced-cycle marginals). "
     "See motion_corrected\spout_position_recovery_cam1\ for the validation montage.")
p=Presentation(DECK); n=0
for s in p.slides:
    for sh in s.shapes:
        if not sh.has_text_frame: continue
        if "SPOUT POSITIONS INCOMPLETE" in sh.text_frame.text:
            # capture original size from first run
            sz=None
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.size: sz=r.font.size; break
                if sz: break
            tf=sh.text_frame; tf.text=NEW
            r0=tf.paragraphs[0].runs[0]
            r0.font.size=sz or Pt(9); r0.font.bold=True
            r0.font.color.rgb=RGBColor(0x1E,0x7A,0x1E)  # dark green = resolved
            n+=1
p.save(DECK)
print(f"updated {n} PS93 8/5 caveat boxes -> resolved recovery note")
