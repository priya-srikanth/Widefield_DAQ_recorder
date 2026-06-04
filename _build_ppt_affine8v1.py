"""Append an 8-point-affine (landmarks v1) re-alignment section to the deck.

Copies E:\\labcams_data\\PS92_94_95.pptx -> PS92_94_95_affine8v1.pptx (original is
NOT modified) and appends a section that is explicit about the transform: each of
the six sessions gets a divider + mean/Allen overlay + cue pre/post/delta + cue
pairwise delta + lick post-by-spout + lick pairwise (delta-position) + cue-vs-lick,
all from the affine8v1 outputs. Sizing matches the existing per-animal slides.

Idempotent: rebuilds from the original each run and includes any session whose
mean-overlay PNG exists, so re-running after 6/3 PS94/PS95 finish yields the full
six-session deck.
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

SRC = r"E:\labcams_data\PS92_94_95.pptx"
DST = r"E:\labcams_data\PS92_94_95_affine8v1.pptx"
TAG = "affine8v1"
D = r"E:\labcams_data"

# (title, date_subtitle, motion_corrected_dir, label, frame-mapping note)
SESSIONS = [
    ("PS94 - 2026-06-01", "full-FOV; raw//2 frame mapping",
     fr"{D}\20260601\PS94_20260601_141614\motion_corrected", "PS94_affine8v1"),
    ("PS95 - 2026-06-01", "full-FOV; raw//2 frame mapping",
     fr"{D}\20260601\PS95_20260601_153653\motion_corrected", "PS95_affine8v1"),
    ("PS92 - 2026-06-02 (rescued)", "ROI crop; relabeled cleanpairs frame-map mapping; functional-channel fix",
     fr"{D}\20260602\PS92\PS92_20260602_151820\illuminated_rescue\motion_corrected", "PS92_0602_affine8v1"),
    ("PS92 - 2026-06-03", "ROI crop; relabeled cleanpairs frame-map mapping",
     fr"{D}\20260603\PS92\PS92_20260603_104008\motion_corrected", "PS92_0603_affine8v1"),
    ("PS94 - 2026-06-03", "ROI crop; relabeled cleanpairs frame-map mapping",
     fr"{D}\20260603\PS94\motion_corrected", "PS94_0603_affine8v1"),
    ("PS95 - 2026-06-03", "ROI crop; relabeled cleanpairs frame-map mapping",
     fr"{D}\20260603\PS95\PS95\PS95_20260603_194442\motion_corrected", "PS95_0603_affine8v1"),
]

TRANSFORM_NOTE = ("8-point AFFINE transform (OB_center/L/R, RSP_base, MOp_L/R, SS_L/R), "
                  "hand-placed landmarks v1; ROI-aware warp to the 540x640 Allen atlas grid. "
                  "Differs from the earlier deck sections, which used a 4-point similarity transform.")

shutil.copy2(SRC, DST)
prs = Presentation(DST)
blank = next(l for l in prs.slide_layouts if l.name == "Blank")


def txt(slide, l, t, w, h, s, size, color, bold=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = s
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)
    return tb


def set_dark_bg(slide, hexcolor="222222"):
    ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    cSld = slide._element.find(f"{ns}cSld")
    bg = etree.SubElement(cSld, f"{ns}bg"); bgPr = etree.SubElement(bg, f"{ns}bgPr")
    sf = etree.SubElement(bgPr, f"{{{a}}}solidFill")
    etree.SubElement(sf, f"{{{a}}}srgbClr").set("val", hexcolor)
    etree.SubElement(bgPr, f"{{{a}}}effectLst")
    cSld.insert(0, bg)


def content(title, subtitle, img, l, t, w):
    s = prs.slides.add_slide(blank)
    txt(s, 0.44, 0.25, 12.44, 0.32, title, 19, "222222", bold=True)
    txt(s, 0.44, 0.64, 12.44, 0.30, subtitle, 8.5, "666666")
    if img and os.path.exists(img):
        s.shapes.add_picture(img, Inches(l), Inches(t), width=Inches(w))
    else:
        txt(s, 0.44, 3.5, 12.44, 0.5, f"[missing image: {img}]", 11, "AA2222")
    txt(s, 0.44, 7.19, 12.44, 0.14, os.path.basename(img) if img else "", 6.5, "666666")
    return s


def divider(title, subtitle):
    d = prs.slides.add_slide(blank); set_dark_bg(d)
    txt(d, 0.64, 2.85, 11.81, 0.9, title, 40, "FFFFFF", bold=True)
    txt(d, 0.69, 3.95, 11.53, 0.8, subtitle, 13, "D8D8D8")


# ---- section intro ----
divider("8-point affine re-alignment (landmarks v1)", TRANSFORM_NOTE)

included = []
for title, datesub, mc, lab in SESSIONS:
    sp = os.path.join(mc, f"spout_trial_averages_{TAG}")
    lk = os.path.join(mc, f"lick_aligned_{TAG}")
    mean_png = os.path.join(sp, f"{lab}_mean_415_470_with_allen_overlay.png")
    if not os.path.exists(mean_png):
        print("skip (no figures yet):", title, lab)
        continue
    included.append(title)
    sub = f"8-pt affine, landmarks v1 - {datesub}"

    divider(title + "  (8-pt affine, v1)", TRANSFORM_NOTE)
    content(f"{title}: mean 415/470 nm with Allen outlines", sub,
            mean_png, 0.10, 1.95, 13.14)
    content(f"{title}: cue-aligned spout averages (1 s pre, 1 s post, post-pre delta)", sub,
            os.path.join(sp, f"{lab}_spout_positions_1s_pre_post_delta_allen_overlay.png"), 4.91, 1.59, 3.51)
    content(f"{title}: cue pairwise delta-position contrasts", sub,
            os.path.join(sp, f"{lab}_pairwise_spout_position_delta_contrasts_allen_overlay.png"), 4.01, 1.59, 5.31)
    content(f"{title}: post-lick maps by spout position (150 ms)", sub,
            os.path.join(lk, f"{lab}_lick_aligned_150ms_post_by_spout.png"), 2.60, 1.40, 8.20)
    content(f"{title}: post-lick pairwise delta-position contrasts", sub,
            os.path.join(lk, f"{lab}_lick_aligned_pairwise_spout_position_contrasts.png"), 4.01, 1.59, 5.31)
    content(f"{title}: cue-aligned vs lick-aligned maps", sub,
            os.path.join(lk, f"{lab}_cue_vs_lick_spout_position_maps.png"), 5.00, 1.59, 3.20)

prs.save(DST)
print(f"saved {DST}  | slides: {len(prs.slides)}  | sessions added: {included}")
