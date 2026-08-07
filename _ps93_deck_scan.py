"""Scan the deck for PS93 2026-08-05 slides and report any red text boxes (stale 4/6 caveat)."""
from pptx import Presentation
from pptx.util import Emu
DECK=r"N:\MICROSCOPE\Priya\Widefield\labcams\PS92_94_95_affine8v1.pptx"
p=Presentation(DECK)
def txt(s):
    return " ".join(t.text for sh in s.shapes if sh.has_text_frame for para in sh.text_frame.paragraphs for t in para.runs)
for i,s in enumerate(p.slides):
    head=txt(s)[:90]
    if "PS93" in head and "2026-08-05" in head:
        print(f"\n=== slide {i}: {head}")
        for j,sh in enumerate(s.shapes):
            if not sh.has_text_frame: continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    col=None
                    try:
                        if r.font.color and r.font.color.type is not None: col=str(r.font.color.rgb)
                    except Exception: col="?"
                    if r.text.strip() and (col in ("FF0000","C00000","FF0000") or "4/6" in r.text or "4 of 6" in r.text or "contaminat" in r.text.lower() or "dead strobe" in r.text.lower()):
                        print(f"   shape{j} col={col} txt={r.text[:120]!r}")
