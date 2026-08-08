"""Build cross-session_aligned.pptx: per-animal sections, dates 6/5-6/8, equivalent
images grouped by family across dates. All in 6/6 CCF (allen_aligned_affine8v1).
Order per animal: mean+Allen -> cue -> lick -> quietnorm lick -> per-session photobleach
-> alignment overlay (vs 6/6 reference) -> motion QC. Cross-date photobleaching once up top.
Skips any missing figure (prints what's missing). Idempotent: rebuilds from scratch.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt

OUT = r"N:\MICROSCOPE\Priya\Widefield\labcams\cross-session_aligned.pptx"
NL = r"N:\MICROSCOPE\Priya\Widefield\labcams"
XINT = r"C:\Github\Widefield_DAQ_recorder\_crossday_intensity_out\crossday_raw_intensity.png"
PBDIR = lambda d: rf"C:\Github\Widefield_DAQ_recorder\_photobleach_out_{d}"

ANIMALS = ["PS92", "PS93", "PS94", "PS95"]
import re, glob
# Auto-discover task-era sessions (6/5 onward) from N: so new days appear with no code edit.
MIN_DATE = "20260605"   # cross-registered era; excludes 6/1-6/4 full-FOV baseline
def _discover():
    sess = {a: {} for a in ANIMALS}
    for dd in sorted(os.listdir(NL)):
        if not re.fullmatch(r"\d{8}", dd) or dd < MIN_DATE:
            continue
        for a in ANIMALS:
            for s in sorted(glob.glob(rf"{NL}\{dd}\{a}_*")):
                if os.path.exists(rf"{s}\motion_corrected\wfield_local_results\frames_average.npy"):
                    sess[a][dd[4:]] = os.path.basename(s); break
    dates = sorted({k for a in ANIMALS for k in sess[a]})
    return sess, dates
SESS, DATES = _discover()
def mc(an, d): return rf"{NL}\2026{d}\{SESS[an][d]}\motion_corrected"
def lab(an, d): return f"{an}_{d}_affine8v1"
# per-(animal,date) family -> figure path
FAMILIES = [
  ("Mean 415/470 + Allen overlay", lambda an,d: rf"{mc(an,d)}\spout_trial_averages_affine8v1\{lab(an,d)}_mean_415_470_with_allen_overlay.png"),
  ("Cue map (2 s post - 2 s pre)", lambda an,d: rf"{mc(an,d)}\spout_trial_averages_affine8v1\{lab(an,d)}_spout_positions_1s_pre_post_delta_shared_scale.png"),
  ("Lick map (150 ms by spout)",   lambda an,d: rf"{mc(an,d)}\lick_aligned_affine8v1\{lab(an,d)}_lick_aligned_150ms_post_by_spout.png"),
  ("Lick map (quiet-normalized)",  lambda an,d: rf"{mc(an,d)}\lick_aligned_affine8v1\{lab(an,d)}_lick_aligned_150ms_post_by_spout_quietnorm.png"),
  ("Per-session photobleaching",   lambda an,d: rf"{PBDIR(d)}\photobleach_{an}_{d}.png"),
  ("Motion QC",                    lambda an,d: rf"{mc(an,d)}\motion_qc\{lab(an,d)}_motion_qc.png"),
]
ALIGN = lambda an: rf"{NL}\xday\{an}_xall\{an}_cross_day_alignment_qc.png"  # per-animal, one fig

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
missing = []
def _title(s, text, size=26):
    tb = s.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.8)); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = True
def title_slide(text, sub=""):
    s = prs.slides.add_slide(BLANK); _title(s, text, 34)
    if sub:
        tb = s.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(12.7), Inches(1.0)); tb.text_frame.word_wrap = True
        tb.text_frame.paragraphs[0].text = sub; tb.text_frame.paragraphs[0].font.size = Pt(16)
def img_slide(text, png, top=1.15, h=6.0):
    s = prs.slides.add_slide(BLANK); _title(s, text)
    if png and os.path.exists(png):
        s.shapes.add_picture(png, Inches(0.4), Inches(top), height=Inches(h)); return True
    tb = s.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(12.5), Inches(1.0))
    tb.text_frame.paragraphs[0].text = "(figure not available)"; tb.text_frame.paragraphs[0].font.size = Pt(18)
    missing.append(text); return False

_dr = f"{int(DATES[0][:2])}/{int(DATES[0][2:])}-{int(DATES[-1][:2])}/{int(DATES[-1][2:])}" if DATES else "n/a"
title_slide(f"Cross-session aligned results ({_dr})",
            "All sessions cross-registered to each animal's 2026-06-06 session (6/6 CCF; "
            "PS92/PS93 v2 landmarks). Grouped by animal; equivalent images across dates. "
            "Auto-discovered task-era sessions (6/5 onward).")
img_slide("Cross-date raw fluorescence intensity (all animals)", XINT, top=1.4, h=5.3)
for an in ANIMALS:
    title_slide(an)
    for fam_name, fam_fn in FAMILIES:
        for d in DATES:
            if d not in SESS[an]:
                continue
            img_slide(f"{an}  -  {fam_name}  -  {int(d[:2])}/{int(d[2:])}", fam_fn(an, d))
        if fam_name == "Per-session photobleaching":  # alignment overlay right after photobleach, before motion QC
            img_slide(f"{an}  -  Alignment to 6/6 reference (vasculature overlay, all days)", ALIGN(an), top=1.3, h=5.6)

# ---- Hemodynamic-correction comparison section (415 vs raw 470 vs corrected 470) ----
def img_centered(text, png, h, aspect, top=1.15):
    s = prs.slides.add_slide(BLANK); _title(s, text)
    if png and os.path.exists(png):
        w = h * aspect; s.shapes.add_picture(png, Inches(max(0.2, (13.333 - w) / 2)), Inches(top), height=Inches(h)); return True
    missing.append(text); return False

CC = os.path.join(NL, "channel_comparison")
title_slide("Hemodynamic correction: 415 nm vs 470 nm vs corrected 470 nm",
            "Example session PS92 6/8. 415 nm (isosbestic) = hemodynamics; 470 nm raw = neural + hemo; "
            "corrected 470 nm = neural (415 regressed out via SVTcorr). Same U_atlas basis; shared color "
            "scale per row; all in 6/6 CCF. (SVGs in labcams\\channel_comparison\\.)")
img_centered("415 vs 470 vs corrected — cue-aligned, pooled (post-cue mean; post-pre delta)",
             os.path.join(CC, "PS92_0608_cue_415_vs_470_vs_corr.png"), 6.1, 13/9.5)
img_centered("415 vs 470 vs corrected — lick-aligned, pooled (150 ms post)",
             os.path.join(CC, "PS92_0608_lick_415_vs_470_vs_corr.png"), 4.4, 13/5.2, top=1.6)
img_centered("415 vs 470 vs corrected — cue-aligned (post-pre) by spout position",
             os.path.join(CC, "PS92_0608_cue_415_vs_470_vs_corr_by_position.png"), 6.4, 13/22, top=1.0)
img_centered("415 vs 470 vs corrected — lick-aligned (150 ms post) by spout position",
             os.path.join(CC, "PS92_0608_lick_415_vs_470_vs_corr_by_position.png"), 6.4, 13/22, top=1.0)

prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides)} slides)")
if missing:
    print(f"\n{len(missing)} missing figures:")
    for m in missing: print("  -", m)
