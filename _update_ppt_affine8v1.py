"""In-place refresh of the affine8v1 deck: swap map pictures for the outline-fixed
versions and append newly-ready session sections, WITHOUT disturbing anything else
(original slides, and user-added reference atlas images are preserved).

- Existing affine8v1 content slides: the one map picture is identified by matching
  its position to the builder layout (the user's reference image sits elsewhere, so
  it is never touched) and replaced in place with the regenerated PNG.
- Sessions not yet in the deck (detected by title) get a fresh divider + 6 slides
  appended at the end.

A timestamped-free backup ``*.bak`` is written before saving.
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

DST = r"N:\MICROSCOPE\Priya\Widefield\labcams\PS92_94_95_affine8v1.pptx"  # canonical deck on MICROSCOPE
TAG = "affine8v1"
D = r"E:\labcams_data"  # figures still read locally from E (SVD/alignment/maps kept there)

# (title, date_subtitle, motion_corrected_dir, label)
SESSIONS = [
    ("PS94 - 2026-06-01", "full-FOV; raw//2 frame mapping",
     fr"{D}\20260601\PS94_20260601_141614\motion_corrected", "PS94_affine8v1"),
    ("PS95 - 2026-06-01", "full-FOV; raw//2 frame mapping",
     fr"{D}\20260601\PS95_20260601_153653\motion_corrected", "PS95_affine8v1"),
    ("PS92 - 2026-06-02 (rescued)", "ROI crop; cleanpairs frame-map mapping; functional-channel fix",
     fr"{D}\20260602\PS92_20260602_151820\illuminated_rescue\motion_corrected", "PS92_0602_affine8v1"),
    ("PS92 - 2026-06-03", "ROI crop; cleanpairs frame-map mapping",
     fr"{D}\20260603\PS92_20260603_104008\motion_corrected", "PS92_0603_affine8v1"),
    ("PS94 - 2026-06-03", "ROI crop; cleanpairs frame-map mapping",
     fr"{D}\20260603\PS94_20260603\motion_corrected", "PS94_0603_affine8v1"),
    ("PS95 - 2026-06-03", "ROI crop; cleanpairs frame-map mapping",
     fr"{D}\20260603\PS95_20260603_194442\motion_corrected", "PS95_0603_affine8v1"),
    ("PS92 - 2026-06-04 (concat)", "ROI crop; cleanpairs frame-map; force-split session concatenated (gap padded)",
     fr"{D}\20260604\PS92_20260604_132934\motion_corrected", "PS92_0604_affine8v1"),
    ("PS94 - 2026-06-04", "ROI crop; cleanpairs frame-map mapping",
     fr"{D}\20260604\PS94_20260604_151516\motion_corrected", "PS94_0604_affine8v1"),
    ("PS95 - 2026-06-04", "ROI crop; cleanpairs frame-map mapping",
     fr"{D}\20260604\PS95_20260604_165712\motion_corrected", "PS95_0604_affine8v1"),
    ("PS92 - 2026-06-05", "trial-triggered; ROI; cleanpairs frame-map; cue 2 s post / 1 s pre",
     fr"{D}\20260605\PS92_20260605_125023\motion_corrected", "PS92_0605_affine8v1"),
    ("PS93 - 2026-06-05 (new animal)", "continuous; ROI; cleanpairs frame-map; cue 2 s post / 1 s pre",
     fr"{D}\20260605\PS93_20260605_174659\motion_corrected", "PS93_0605_affine8v1"),
    ("PS94 - 2026-06-05", "continuous; ROI; cleanpairs frame-map; cue 2 s post / 1 s pre",
     fr"{D}\20260605\PS94_20260605_142009\motion_corrected", "PS94_0605_affine8v1"),
    ("PS95 - 2026-06-05", "continuous; ROI; cleanpairs frame-map; cue 2 s post / 1 s pre",
     fr"{D}\20260605\PS95_20260605_163102\motion_corrected", "PS95_0605_affine8v1"),
    ("PS92 - 2026-06-06", "continuous; ROI; cleanpairs frame-map; cue 2 s post / 2 s pre",
     fr"{D}\20260606\PS92_20260606_122451\motion_corrected", "PS92_0606_affine8v1"),
    ("PS93 - 2026-06-06", "continuous; ROI; cleanpairs frame-map; cue 2 s post / 2 s pre",
     fr"{D}\20260606\PS93_20260606_180117\motion_corrected", "PS93_0606_affine8v1"),
    ("PS94 - 2026-06-06", "continuous; ROI; cleanpairs frame-map; cue 2 s post / 2 s pre",
     fr"{D}\20260606\PS94_20260606_140854\motion_corrected", "PS94_0606_affine8v1"),
    ("PS95 - 2026-06-06", "continuous; ROI; cleanpairs frame-map; cue 2 s post / 2 s pre",
     fr"{D}\20260606\PS95_20260606_160806\motion_corrected", "PS95_0606_affine8v1"),
    ("PS92 - 2026-06-07", "FIXED motion; cross-registered to 6/6 (6/6 CCF); cue 2 s post / 2 s pre",
     fr"{D}\20260607\PS92_20260607_121538\motion_corrected", "PS92_0607_affine8v1"),
    ("PS93 - 2026-06-07", "FIXED motion; cross-registered to 6/6 (6/6 CCF); cue 2 s post / 2 s pre",
     fr"{D}\20260607\PS93_20260607_174844\motion_corrected", "PS93_0607_affine8v1"),
    ("PS94 - 2026-06-07", "FIXED motion; cross-registered to 6/6 (6/6 CCF); cue 2 s post / 2 s pre",
     fr"{D}\20260607\PS94_20260607_140731\motion_corrected", "PS94_0607_affine8v1"),
    ("PS95 - 2026-06-07", "FIXED motion; cross-registered to 6/6 (6/6 CCF); cue 2 s post / 2 s pre",
     fr"{D}\20260607\PS95_20260607_155000\motion_corrected", "PS95_0607_affine8v1"),
    ("PS92 - 2026-06-08", "FIXED motion; cross-registered to 6/6 (6/6 CCF); cue 2 s post / 2 s pre",
     fr"{D}\20260608\PS92_20260608_133759\motion_corrected", "PS92_0608_affine8v1"),
    ("PS93 - 2026-06-08", "FIXED motion; cross-registered to 6/6 (6/6 CCF); cue 2 s post / 2 s pre",
     fr"{D}\20260608\PS93_20260608_195203\motion_corrected", "PS93_0608_affine8v1"),
    ("PS94 - 2026-06-08", "FIXED motion; cross-registered to 6/6 (6/6 CCF); cue 2 s post / 2 s pre",
     fr"{D}\20260608\PS94_20260608_153651\motion_corrected", "PS94_0608_affine8v1"),
    ("PS95 - 2026-06-08", "FIXED motion; cross-registered to 6/6 (6/6 CCF); cue 2 s post / 2 s pre",
     fr"{D}\20260608\PS95_20260608_180943\motion_corrected", "PS95_0608_affine8v1"),
]
TRANSFORM_NOTE = ("8-point AFFINE transform (OB_center/L/R, RSP_base, MOp_L/R, SS_L/R), "
                  "hand-placed landmarks v1; ROI-aware warp to the 540x640 Allen atlas grid. "
                  "Differs from the earlier deck sections, which used a 4-point similarity transform.")

# per-figure: (title_suffix, subdir, filename_suffix, left, top, width)
FIGS = [
    (": mean 415/470 nm with Allen outlines", "spout", "_mean_415_470_with_allen_overlay.png", 0.10, 1.95, 13.14),
    (": cue-aligned spout averages (pre / post / post-pre delta, shared scale; window per panel labels)", "spout", "_spout_positions_1s_pre_post_delta_shared_scale.png", 4.91, 1.59, 3.51),
    (": cue pairwise delta-position contrasts", "spout", "_pairwise_spout_position_delta_contrasts_allen_overlay.png", 4.01, 1.59, 5.31),
    (": post-lick maps by spout position (150 ms)", "lick", "_lick_aligned_150ms_post_by_spout.png", 2.60, 1.40, 8.20),
    (": post-lick pairwise delta-position contrasts", "lick", "_lick_aligned_pairwise_spout_position_contrasts.png", 4.01, 1.59, 5.31),
    (": cue-aligned vs lick-aligned maps", "lick", "_cue_vs_lick_spout_position_maps.png", 5.00, 1.59, 3.20),
]


def fig_path(mc, lab, subdir, suffix):
    base = os.path.join(mc, f"spout_trial_averages_{TAG}" if subdir == "spout" else f"lick_aligned_{TAG}")
    return os.path.join(base, f"{lab}{suffix}")


def slide_title(s):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip()
    return ""


# title -> (png, l, t, w) for every figure of every session that has figures
want = {}
sessions_present_figs = []
for title, datesub, mc, lab in SESSIONS:
    mean = fig_path(mc, lab, "spout", "_mean_415_470_with_allen_overlay.png")
    if not os.path.exists(mean):
        continue
    sessions_present_figs.append((title, datesub, mc, lab))
    for suffix_title, subdir, fsuffix, l, t, w in FIGS:
        want[title + suffix_title] = (fig_path(mc, lab, subdir, fsuffix), l, t, w)

shutil.copy2(DST, DST + ".bak")
prs = Presentation(DST)

# ---- 1) swap map pictures on existing content slides ----
existing_titles = set()
swapped = 0
for s in prs.slides:
    ttl = slide_title(s)
    existing_titles.add(ttl)
    if ttl not in want:
        continue
    png, l, t, w = want[ttl]
    pics = [sh for sh in s.shapes if sh.shape_type == 13]
    if not pics:
        continue
    # the map picture is the one whose position matches the builder layout;
    # the user's reference image sits elsewhere and is left untouched.
    target = min(pics, key=lambda sh: abs(Emu(sh.left).inches - l) + abs(Emu(sh.top).inches - t))
    if abs(Emu(target.left).inches - l) + abs(Emu(target.top).inches - t) > 0.2:
        print(f"  WARN: no map pic match on '{ttl}' (skipped, nothing removed)")
        continue
    L, T = target.left, target.top
    target._element.getparent().remove(target._element)
    if os.path.exists(png):
        s.shapes.add_picture(png, L, T, width=Inches(w))
        swapped += 1
    else:
        print(f"  WARN: regenerated PNG missing: {png}")

# ---- 2) append sections for sessions not yet in the deck ----
blank = next(layout for layout in prs.slide_layouts if layout.name == "Blank")


def txt(slide, l, t, w, h, s, size, color, bold=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = s
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)


def set_dark_bg(slide, hexcolor="222222"):
    ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    cSld = slide._element.find(f"{ns}cSld")
    bg = etree.SubElement(cSld, f"{ns}bg"); bgPr = etree.SubElement(bg, f"{ns}bgPr")
    sf = etree.SubElement(bgPr, f"{{{a}}}solidFill")
    etree.SubElement(sf, f"{{{a}}}srgbClr").set("val", hexcolor)
    etree.SubElement(bgPr, f"{{{a}}}effectLst")
    cSld.insert(0, bg)


def content(title, subtitle, img, l, t, w):
    s = prs.slides.add_slide(blank)
    txt(s, 0.44, 0.25, 12.44, 0.32, title, 19, "222222", bold=True)
    txt(s, 0.44, 0.64, 12.44, 0.30, subtitle, 8.5, "666666")
    if img and os.path.exists(img):
        s.shapes.add_picture(img, Inches(l), Inches(t), width=Inches(w))
    else:
        txt(s, 0.44, 3.5, 12.44, 0.5, f"[missing image: {img}]", 11, "AA2222")
    txt(s, 0.44, 7.19, 12.44, 0.14, os.path.basename(img) if img else "", 6.5, "666666")


def divider(title, subtitle):
    d = prs.slides.add_slide(blank); set_dark_bg(d)
    txt(d, 0.64, 2.85, 11.81, 0.9, title, 40, "FFFFFF", bold=True)
    txt(d, 0.69, 3.95, 11.53, 0.8, subtitle, 13, "D8D8D8")


added = []
for title, datesub, mc, lab in sessions_present_figs:
    if (title + FIGS[0][0]) in existing_titles:
        continue  # already in the deck (pictures were swapped above)
    added.append(title)
    sub = f"8-pt affine, landmarks v1 - {datesub}"
    divider(title + "  (8-pt affine, v1)", TRANSFORM_NOTE)
    for suffix_title, subdir, fsuffix, l, t, w in FIGS:
        content(title + suffix_title, sub, fig_path(mc, lab, subdir, fsuffix), l, t, w)

# ---- 3) Motion-correction QC section (idempotent via title guards) ----
present = {slide_title(s) for s in prs.slides}
qc_div = "Motion-correction QC"
qc_sessions = [(t, mc, lab) for t, _, mc, lab in sessions_present_figs
               if os.path.exists(os.path.join(mc, "motion_qc", f"{lab}_motion_qc.png"))]
qc_added = []
if qc_sessions and qc_div not in present:
    divider(qc_div, "Per-frame rigid shifts + histogram, mean-image sharpness (raw vs corrected), "
                    "and corrected temporal-std (residual motion). Median shift should be sub-pixel.")
_qc_slides = {slide_title(s): s for s in prs.slides}
for title, mc, lab in qc_sessions:
    ct = f"{title}: motion-correction QC"
    qcpng = os.path.join(mc, "motion_qc", f"{lab}_motion_qc.png")
    if ct in _qc_slides:   # refresh the QC image in place (e.g. after re-doing motion correction)
        sl = _qc_slides[ct]
        for sh in [sh for sh in sl.shapes if sh.shape_type == 13]:
            sh._element.getparent().remove(sh._element)
        sl.shapes.add_picture(qcpng, Inches(0.15), Inches(1.35), width=Inches(12.9))
        qc_added.append(title + "(refresh)")
    else:
        content(ct, "Shift traces / magnitude histogram / sharpness / residual-motion std",
                qcpng, 0.15, 1.35, 12.9)
        qc_added.append(title)

# ---- 4) Remove the stale "LED drift across sessions" slide (6/1-6/3 only;
#         superseded by the per-day photobleach slides + the cross-day intensity slide). ----
pb_removed = []
_stale_title = "Photobleaching / LED drift across sessions"
_sldIdLst = prs.slides._sldIdLst
while True:
    _ids = list(_sldIdLst); _hit = False
    for _idx, _sl in enumerate(prs.slides):
        if slide_title(_sl) == _stale_title:
            _sldIdLst.remove(_ids[_idx]); pb_removed.append(_stale_title); _hit = True; break
    if not _hit:
        break

# ---- 4a) 2026-06-05 photobleaching: continuous vs trial-triggered (idempotent) ----
present = {slide_title(s) for s in prs.slides}
pb0605_dir = r"C:\Github\Widefield_DAQ_recorder\_photobleach_out_0605"
pb0605_sum = os.path.join(pb0605_dir, "photobleach_0605_continuous_vs_trial.png")
pb0605_title = "Photobleaching 2026-06-05: continuous vs trial-triggered"
pb0605_added = []
if os.path.exists(pb0605_sum) and pb0605_title not in present:
    content(pb0605_title,
            "Functional 470 nm drift: trial-triggered (PS92) +0.04%/min vs continuous (PS93/94/95) "
            "mean -0.02%/min -> continuous only modestly worse (~0.06%/min). Larger 415 isosbestic "
            "decline (esp. PS93, longest) = violet-LED drift, removed by the 0.1 Hz hemo highpass.",
            pb0605_sum, 0.15, 1.7, 12.9)
    pb0605_added.append("summary")
    for an in ("PS92", "PS93", "PS94", "PS95"):
        per = os.path.join(pb0605_dir, f"photobleach_{an}_0605.png")
        pt = f"Photobleaching 2026-06-05: {an} per-channel trend"
        if os.path.exists(per) and pt not in present:
            content(pt, "415 isosbestic vs 470 functional ROI-median intensity (binned + linear fit)",
                    per, 0.6, 1.7, 12.0)
            pb0605_added.append(an)

# ---- 4a2) 2026-06-06 photobleaching (idempotent) ----
present = {slide_title(s) for s in prs.slides}
pb0606_dir = r"C:\Github\Widefield_DAQ_recorder\_photobleach_out_0606"
pb0606_sum = os.path.join(pb0606_dir, "photobleach_SUMMARY.png")
pb0606_title = "Photobleaching 2026-06-06 (4 sessions)"
pb0606_added = []
if os.path.exists(pb0606_sum) and pb0606_title not in present:
    content(pb0606_title,
            "415 isosbestic vs 470 functional ROI-median trends + per-session %drift. "
            "470 stays within a few % (PS92 -2.9, PS93 -1.5, PS94 +0.4, PS95 -6.1); "
            "larger 415 decline = violet-LED drift (removed by the 0.1 Hz hemo highpass).",
            pb0606_sum, 0.15, 1.7, 12.9)
    pb0606_added.append("summary")
    for an in ("PS92", "PS93", "PS94", "PS95"):
        per = os.path.join(pb0606_dir, f"photobleach_{an}_0606.png")
        pt = f"Photobleaching 2026-06-06: {an} per-channel trend"
        if os.path.exists(per) and pt not in present:
            content(pt, "415 vs 470 ROI-median intensity (binned + linear fit)", per, 0.6, 1.7, 12.0)
            pb0606_added.append(an)

# ---- 4a3) 2026-06-07 photobleaching (idempotent) ----
present = {slide_title(s) for s in prs.slides}
pb0607_dir = r"C:\Github\Widefield_DAQ_recorder\_photobleach_out_0607"
pb0607_sum = os.path.join(pb0607_dir, "photobleach_SUMMARY.png")
pb0607_title = "Photobleaching 2026-06-07 (4 sessions)"
if os.path.exists(pb0607_sum) and pb0607_title not in present:
    content(pb0607_title,
            "415 isosbestic vs 470 functional ROI-median trends + per-session %drift. "
            "470 stays within a few % (PS92 -2.1, PS93 +1.3, PS94 +1.3, PS95 -4.1); "
            "larger 415 decline = violet-LED drift (removed by the 0.1 Hz hemo highpass).",
            pb0607_sum, 0.15, 1.7, 12.9)
    for an in ("PS92", "PS93", "PS94", "PS95"):
        per = os.path.join(pb0607_dir, f"photobleach_{an}_0607.png")
        pt = f"Photobleaching 2026-06-07: {an} per-channel trend"
        if os.path.exists(per) and pt not in present:
            content(pt, "415 vs 470 ROI-median intensity (binned + linear fit)", per, 0.6, 1.7, 12.0)

# ---- 4a3b) 2026-06-08 photobleaching (idempotent) ----
present = {slide_title(s) for s in prs.slides}
pb0608_dir = r"C:\Github\Widefield_DAQ_recorder\_photobleach_out_0608"
pb0608_sum = os.path.join(pb0608_dir, "photobleach_SUMMARY.png")
pb0608_title = "Photobleaching 2026-06-08 (4 sessions)"
if os.path.exists(pb0608_sum) and pb0608_title not in present:
    content(pb0608_title,
            "415 isosbestic vs 470 functional ROI-median trends + per-session %drift. "
            "470 stays within a few % (PS92 -4.8, PS93 +4.8, PS94 -0.9, PS95 -4.9); "
            "larger 415 decline = violet-LED drift (removed by the 0.1 Hz hemo highpass).",
            pb0608_sum, 0.15, 1.7, 12.9)
    for an in ("PS92", "PS93", "PS94", "PS95"):
        per = os.path.join(pb0608_dir, f"photobleach_{an}_0608.png")
        pt = f"Photobleaching 2026-06-08: {an} per-channel trend"
        if os.path.exists(per) and pt not in present:
            content(pt, "415 vs 470 ROI-median intensity (binned + linear fit)", per, 0.6, 1.7, 12.0)

# ---- 4a4) Cross-day RAW fluorescence intensity (idempotent) ----
present = {slide_title(s) for s in prs.slides}
xint_png = r"C:\Github\Widefield_DAQ_recorder\_crossday_intensity_out\crossday_raw_intensity.png"
xint_title = "Cross-day raw fluorescence intensity (per animal)"
xint_caption = ("Brain-ROI median raw counts (415 + 470) per animal across days, from each "
                "session's frames_average. CAVEAT: LED power is manually titrated day-to-day, so "
                "a trend may reflect the LED setting, not photobleaching. No monotonic decline seen.")
if os.path.exists(xint_png):
    _xint = {slide_title(s): s for s in prs.slides}.get(xint_title)
    if _xint is not None:   # refresh the figure in place (latest days added)
        for sh in [sh for sh in _xint.shapes if sh.shape_type == 13]:
            sh._element.getparent().remove(sh._element)
        _xint.shapes.add_picture(xint_png, Inches(0.4), Inches(1.7), width=Inches(12.5))
    else:
        content(xint_title, xint_caption, xint_png, 0.4, 1.7, 12.5)

# ---- 4c) Cross-day vasculature alignment QC (idempotent) ----
present = {slide_title(s) for s in prs.slides}
xday_root = r"N:\MICROSCOPE\Priya\Widefield\xday"
xday_div = "Cross-day alignment (vasculature)"
xday_added = []
xday_imgs = []
if os.path.isdir(xday_root):
    for an in sorted(os.listdir(xday_root)):
        qc = os.path.join(xday_root, an, f"{an}_cross_day_alignment_qc.png")
        if os.path.exists(qc):
            xday_imgs.append((an, qc))
if xday_imgs and xday_div not in present:
    divider(xday_div, "Same-animal across-day registration of the motion-corrected mean 470 nm "
                      "vasculature to a reference session (ROI days; 540x640 CCF frame). "
                      "NCC = masked vessel correlation vs reference (1.0 = reference; higher = better).")
_xday_slides = {slide_title(s): s for s in prs.slides}
for an, qc in xday_imgs:
    ct = f"{an}: cross-day vasculature alignment QC"
    if ct in _xday_slides:   # refresh the embedded image (e.g. greedy -> reference-native)
        sl = _xday_slides[ct]
        for sh in [sh for sh in sl.shapes if sh.shape_type == 13]:
            sh._element.getparent().remove(sh._element)
        sl.shapes.add_picture(qc, Inches(0.5), Inches(1.4), width=Inches(12.4))
        xday_added.append(an + "(refresh)")
    else:
        content(ct, "red/green vessel overlay per day + masked NCC vs reference (reference-native)", qc, 0.5, 1.4, 12.4)
        xday_added.append(an)

# ---- 4b) Quiet-normalized lick activity (lands before QC; QC is moved to end below) ----
NLAB = r"N:\MICROSCOPE\Priya\Widefield\labcams"
present = {slide_title(s) for s in prs.slides}
ql_div = "Quiet-normalized lick activity"
ql_added = []
ql_sessions = []
for title, datesub, mc, lab in SESSIONS:
    n_mc = mc.replace(r"E:\labcams_data", NLAB)
    qn = os.path.join(n_mc, "lick_aligned_affine8v1", f"{lab}_lick_aligned_150ms_post_by_spout_quietnorm.png")
    if os.path.exists(qn):
        ql_sessions.append((title, qn))
if ql_sessions and ql_div not in present:
    divider(ql_div, "150 ms post-lick by spout position MINUS the mean quiet-period (not-running / "
                    "not-licking) baseline = lick-evoked activity relative to the quiet state. "
                    "Quiet-period thresholds are provisional (tune later).")
for title, qn in ql_sessions:
    ct = f"{title}: lick-evoked vs quiet baseline (150 ms post-lick)"
    if ct not in present:
        content(ct, "post-lick by spout position minus quiet-period baseline (quiet-normalized)", qn, 2.60, 1.40, 8.20)
        ql_added.append(title)

# ---- 5) Move all QC / diagnostic slides to the very end (after activity maps) ----
def _is_qc(t: str) -> bool:
    t = t.lower()
    return ("motion-correction qc" in t) or ("photobleach" in t) or ("frame alignment qc" in t)

lst = prs.slides._sldIdLst
children = list(lst)
titles = [slide_title(s) for s in prs.slides]
moved = 0
for el, t in zip(children, titles):
    if _is_qc(t):
        lst.remove(el)
        lst.append(el)   # append in original relative order -> stable, idempotent
        moved += 1

prs.save(DST)
print(f"swapped {swapped} map pictures; appended sections: {added}; QC slides: {qc_added}; "
      f"removed-stale: {pb_removed}; quiet-lick: {ql_added}; moved-to-end: {moved}; slides now {len(prs.slides)}")
print(f"backup at {DST}.bak")
