"""Update PS93 8/5 caveat boxes: recovery is now human-verified (279/279 confirmed, 0 changes)."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
DECK=r"N:\MICROSCOPE\Priya\Widefield\labcams\PS92_94_95_affine8v1.pptx"
NEW=("✓ All 6 spout positions RECOVERED from cam1 head-on video and HUMAN-VERIFIED "
     "(279/279 ambiguous trials confirmed in the review GUI; 0 corrections). DAQ strobe bit1 "
     "was dead 8/5 AND this session's behavior log is empty, so positions were inferred from "
     "the spout's x-position in each DAQ-sync-aligned cam1 frame (discrete clusters per fixed "
     "detent), then spot-checked trial-by-trial. See motion_corrected\spout_position_recovery_cam1\.")
p=Presentation(DECK); n=0
for s in p.slides:
    for sh in s.shapes:
        if sh.has_text_frame and "RECOVERED from cam1" in sh.text_frame.text:
            sz=None
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.size: sz=r.font.size; break
                if sz: break
            tf=sh.text_frame; tf.text=NEW; r0=tf.paragraphs[0].runs[0]
            r0.font.size=sz or Pt(9); r0.font.bold=True; r0.font.color.rgb=RGBColor(0x1E,0x7A,0x1E); n+=1
p.save(DECK); print(f"updated {n} PS93 8/5 caveat boxes -> human-verified note")
